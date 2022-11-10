
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
#library to work colors of a parragraph in pptx
from src.utils.template_utils import pptx2pdf, textStyle, elements2change
from src.config import temPath, certificate_template

# use the template to generate the certifcate

def create(data):
    #opening template
    
    template = Presentation(certificate_template)
    content = elements2change(
        nombre = data['nombre'],
        puesto = data['puesto'] if 'puesto' in data else 'Asistente',
        lugar = data['lugar'] if 'lugar' in data else False,
        title = data['titulo'] if 'titulo' in data else False,
        year = data['anio'] if 'anio' in data else 21
    )
    template = writeParagraphs(content,template)
    filename = data['id_constancia']
    template.save(f'{temPath}{filename}.pptx')
    #if pptx2pdf(filename, data['id_constancia']):
    #    return f'{filename}.pdf'
    #else:
    #    return False
    return filename

def writeParagraphs(content, template):
    for idx in content['paragraphs']:
        certificate = template.slides[0].shapes[0].shapes[idx].text_frame
        certificate.clear() 
        certificate.alignment = PP_ALIGN.CENTER
        certificate = certificate.paragraphs[0]
        if content['isFragmented'][idx]:
            references = content['text_fragmented'][idx]
            for fragment in references['paragraphs']:
                paragraph_fragment = textStyle(
                    insert = certificate.add_run(), 
                    isBold = references['bold'][fragment],
                    size = references['size'],
                    fontName = references['font'],
                    color = references['color'],
                )
                paragraph_fragment.text = content['text_fragmented'][idx]['paragraphs'][fragment]
        else:
            paragraph = textStyle(
                insert = certificate.add_run(), 
                isBold = content['bold'],
                size = content['size'],
                fontName = content['font'],
                color = content['color'],
            )
            paragraph.text = content['paragraphs'][idx]
            
    return template

    




