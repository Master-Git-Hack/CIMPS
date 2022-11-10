"""File to handle Certificates"""
from os import uname
from os.path import join
from subprocess import call

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

from .config import Config as _Config
from .Drive import Drive


class Certificates(_Config):
    """class of Certifcates to handle creation, export and protection of certificates"""

    _template = Presentation()
    _paragraph = None
    _filename = None
    _drive = Drive()
    _local_command = dict(Linux="libreoffice", Darwin="soffice")
    _schema = "assistance"

    def __init__(self, schema: str = None) -> None:

        if schema is not None:
            self._schema = schema
        self._template = Presentation(join(self.paths.schemas, f"{self._schema}.pptx"))

    def create(self, data: dict) -> str:
        """function to create the certificates
        Args:
            data (dict) : data to handle with
        Returns:
            filename (str) : filename without extension
        """
        _data = self._elements_to_change(
            nombre=data["nombre"],
            puesto=data["puesto"] if "puesto" in data else "Asistente",
            lugar=data["lugar"] or False,
            title=data["titulo"] or False,
        )
        self._write_paragraphs(_data)
        self._filename = data["id_constancia"]
        self._template.save(f"{join(self.paths.temp,self._filename)}.pptx")
        return self._filename

    def export_as_pdf(
        self, email: str, as_local: bool = False, with_drive: bool = False
    ) -> bool:
        """export current file as pdf
        to export with drive you require to pass the email
        """
        if not with_drive:
            as_local = True
        if email is not None and with_drive:
            self._drive.upload(self._filename)

        if as_local:
            call(
                [
                    self._local_command[uname().sysname],
                    "--headless",
                    "--invisible",
                    "--convert-to",
                    "pdf:writer_pdf_Export",
                    "--outdir",
                    self.paths.temp,
                    f"{self._filename}.pptx",
                ]
            )
            return True
        return False

    def protect_pdf(self, email: str) -> bool:
        """protect pdf with email

        Args:
            email (str): _description_

        Returns:
            bool: _description_
        """
        call(["mv", f"{self._filename}.pdf", f"{self._filename}-tmp.pdf"])

        call(
            [
                "pdftk",
                f"{self._filename}-tmp.pdf",
                "cat",
                "output",
                f"{self._filename}.pdf",
                "owner_pw",
                f"""\
{self.ciudad_sede}_{self.estado_sede} \
{self.current_edition} \
{self.cimps_inicio}_@_{self.cimps_termino} \
{self.public_password_for_file} \
{email}""",
                "allow",
                "printing",
            ]
        )

        call(["rm", f"{self._filename}-tmp.pdf", f"{self._filename}.pptx"])
        return True

    def clean_path(self) -> None:
        """function to clean the path of the certificates"""
        call(["rm", "-rf", self.paths.temp])

    def _text_style(
        self,
        insert: Presentation,
        size: int = 10,
        is_bold: bool = False,
        font_name: str = "Montserrat",
        color: RGBColor = RGBColor(0, 0, 0),
    ):
        style = insert.font
        style.bold = is_bold
        style.size = Pt(size)
        style.name = font_name
        style.color.rgb = color
        return insert

    def _elements_to_change(
        self,
        nombre: str = "John Doe",
        puesto: str = "Asistente",
        lugar: bool or str = False,
        title: bool or str = False,
    ) -> dict:
        """function to change the elements of the template
        Args:
            nombre (str, optional): name of the person. Defaults to 'John Doe'.
            puesto (str, optional): position of the person. Defaults to 'Asistente'.
            lugar (bool, optional): place of the event. Defaults to False.
            title (bool, optional): title of the event. Defaults to False.

        Returns:
            dict: dictionary with the elements to change
        """
        return dict(
            paragraphs={
                2: f"{nombre.upper()}",
                3: "",
                4: f"{title.upper() if title is not False else puesto.upper() }",
            },
            font="Calibri",
            color=RGBColor(136, 0, 0),
            bold=True,
            size=32 if len(nombre) <= 30 else 24,
            isFragmented={2: False, 3: True, 4: False},
            text_fragmented={
                3: dict(
                    paragraphs={
                        1: "Por su valiosa participación "
                        + ("" if not title else "como:"),
                        2: f" {puesto.upper()} " if title else "",
                        3: f"en el Congreso Internacional CIMPS 20{self.year} "
                        if not title
                        else "",
                        4: f'{"en el " if lugar else "como: "}',
                        5: f'{lugar or ""}',
                    },
                    font="Calibri",
                    bold={1: False, 2: True, 3: False, 4: False, 5: True},
                    size=16,
                    color=RGBColor(0, 0, 0),
                )
            },
        )

    def _write_paragraphs(self, data: dict) -> None:
        if self._schema == "commite":
            _certificate = self._template.slides[0].shapes[2].text_frame
            _certificate.clear()
            _certificate.alignment = PP_ALIGN.CENTER
            _certificate = _certificate.paragraphs[0]
            _paragraph = self._text_style(
                insert=_certificate.add_run(),
                is_bold=True,
                size=20,
                font_name="Calibri",
                color=RGBColor(0, 0, 0),
            )
            _paragraph.text = data["Nombre"]

        if self._schema == "assistance":
            for idx in data["paragraphs"]:
                _certificate = self._template.slides[0].shapes[0].shapes[idx].text_frame
                _certificate.clear()
                _certificate.alignment = PP_ALIGN.CENTER
                _certificate = _certificate.paragraphs[0]
                if data["isFragmented"][idx]:
                    references = data["text_fragmented"][idx]
                    for fragment in references["paragraphs"]:
                        paragraph_fragment = self._text_style(
                            insert=_certificate.add_run(),
                            is_bold=True,
                            size=20,
                            font_name="Calibri",
                            color=RGBColor(0, 0, 0),
                        )
                        paragraph_fragment.text = data["text_fragmented"][idx][
                            "paragraphs"
                        ][fragment]
                else:
                    paragraph = self._text_style(
                        insert=_certificate.add_run(),
                        is_bold=references["bold"][fragment],
                        size=references["size"],
                        font_name=references["font"],
                        color=references["color"],
                    )
                    paragraph.text = data["paragraphs"][idx]
