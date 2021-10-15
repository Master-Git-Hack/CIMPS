from credenciales.bd_access import cnxn
from credenciales.google_services import service, MediaFileUpload
from random import randrange as rg
from os import getcwd
import subprocess as sp
#library to work with pptx
from pptx import Presentation
#library to work colors of a parragraph in pptx
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

import smtplib
import os
from credenciales.bd_access import cnxn
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.base import MIMEBase
from email import encoders


#eventos sin usuarios registrados
def consulta():
    if cnxn!=False:
        with cnxn.cursor() as cursor:
            query="""
SELECT c.id AS constancia, ae.id_taller AS id_evento, 
	'For <hisher> attendance at the workshop entitled as:' AS tipo, ae.nombretaller AS taller, uc.nombre AS nombre,
    uc.email, uc.genero,c.id_tipo_constancia
FROM constancia c
	INNER JOIN tipo_constancia t ON c.id_tipo_constancia = t.id
    INNER JOIN (
		SELECT id, nombre, email, genero, CASE id_tipo_usuario WHEN 2 THEN 1 WHEN 3 THEN 4 WHEN 4 THEN 5 WHEN 6 THEN 9 END AS id_tipo_constancia
        FROM usuario_constancia
        WHERE id_tipo_usuario = 6
			AND id_estado = 1 AND id_emision = 9
            AND id NOT IN (
				SELECT user_id
				FROM ingsofti_CIMPS3.users_groups
				WHERE group_id = 2
			)
    )uc ON c.id = uc.id AND c.id_tipo_constancia = uc.id_tipo_constancia
    INNER JOIN (
			SELECT u.id, name, email, e.id AS id_taller, nombre_interno AS nombretaller
		FROM users u
			INNER JOIN asistentes_por_fecha af ON u.id = id_user
			INNER JOIN eventos e ON e.id = id_evento
		WHERE emision = 2020
			AND id_tipo_evento = 2
	)ae ON ae.id = c.id AND ae.id_taller = id_evento
WHERE c.id_estado = 1 AND c.id_emision = 9
ORDER BY constancia, id_evento;


            """
            cursor.execute(query)
            rows = cursor.fetchall()
            data=[]
            for row in rows:
#SELECT o2, e.id_locacion, uc.email
                data.append({
                    'id_constancia':row[0],
                    'id_evento':row[1],
                    'Tipo':row[2],
                    'Taller':row[3],
                    'Nombre':row[4],
                    'Email':row[5],
                    'Genero':row[6],
                    'id_tipo_constancia':row[7]
                })
            return data

def update_records(datos,filename,url):
    if cnxn!=False:
        with cnxn.cursor() as cursor:
            query=f"""
                UPDATE 
                    constancia 
                SET 
                    nombre_archivo='{filename}.pdf',
                    url='{url}',
                    generado=1 
                WHERE 
                    id={datos['id_constancia']}
                    AND
                    id_tipo_constancia={datos['id_tipo_constancia']}
            """
            cursor.execute(query)
            cnxn.commit()


def upload_file(filename):
        folder_id = '1ihRYqMT5Pk5VgycLYY8DrZWm_mkSlVXe'
        myfile=f'{filename}.pdf'
        file_metadata = {
            'name': myfile,
            'parents': [folder_id]
        }
        
        media = MediaFileUpload(f'./archivos/{myfile}', mimetype="application/pdf")
        response = service.files().create(
            body=file_metadata,
            media_body=media,
            fields="id"
        ).execute()

        if response:
            url= f"https://drive.google.com/uc?export=download&id={response.get('id')}"
            return url

def template(info):
    #opening template
    document = Presentation('./template/template_taller.pptx')
    #read specific text frame to edit
    #to get access to other frames changes index of shapes and add .text property after .text_frame
    #shape 3 = In recognition ...
    shape_for_recognition = document.slides[0].shapes[4].text_frame
    shape_for_recognition.clear()
    access_for_recognition=shape_for_recognition.paragraphs[0]
    access_for_recognition.alignment = PP_ALIGN.CENTER
    #giving information to paragraph
    texto1=access_for_recognition.add_run()
    texto1.text = info['Tipo'].replace("<hisher>",info['Genero'][0:3])
    
    #Giving format to name
    font = texto1.font
    font.name = "Helvetica"
    font.size = Pt(15)
    font.color.rgb = RGBColor(255, 255, 255)

    #Shape 1 Name of author
    shape_for_name = document.slides[0].shapes[1].text_frame
    shape_for_name.clear()
    access_for_name=shape_for_name.paragraphs[0]
    access_for_name.alignment = PP_ALIGN.CENTER
    #giving information to paragraph
    name=access_for_name.add_run()
    str_tmp =info['Nombre']
    name.text= ' '.join(str_tmp[1:])
    dcs=['Mr. ','Alum. ','Dr. ','Miss. ','Prof. ','Comp. ','Eng. ']  
    for dc in dcs:
        if dc in str_tmp:
            str_tmp=str_tmp.split()
            str_tmp= ' '.join(str_tmp[1:])
            name.text=str_tmp.upper()
        else:
            name.text=str_tmp.upper()
    #Giving format to name
    font = name.font
    font.name = "Helvetica"
    font.size = Pt(20)
    font.color.rgb = RGBColor(255, 191, 0)

    #Shape 10 name of the exposition
    shape_for_expo = document.slides[0].shapes[10].text_frame
    shape_for_expo.clear()
    access_for_expo=shape_for_expo.paragraphs[0]
    access_for_expo.alignment = PP_ALIGN.CENTER
    #giving information to paragraph
    texto4=access_for_expo.add_run()
    texto4.text = info['Taller'].upper()
    
    #Giving format to name
    font = texto4.font
    font.name = "Helvetica"
    font.size = Pt(20)
    font.color.rgb = RGBColor(255, 191, 0)

    #Saving file
    filename = f"{info['id_constancia']}_{info['Email']}_Taller_{info['Taller'][:11]}"
    document.save(f"./archivos/{filename}-tmp.pptx")
    return filename

def pptx2pdf(filename,email):
    sp.call(['libreoffice', '--headless', '--convert-to', 'pdf', f'./archivos/{filename}-tmp.pptx'])
    #sp.call(['convert',f'./archivos/{filename}.pptx',f'./archivos/{filename}.pdf'])
   # sp.call(['convert',f'./archivos/{filename}.pptx',f'./archivos/{filename}-tmp.pdf'])
    sp.call(['mv',f'{filename}-tmp.pdf','./archivos'])
    sp.call(['pdftk',f'./archivos/{filename}-tmp.pdf','output',f'./archivos/{filename}.pdf','owner_pw',f'#CiMPS_@2020_{email}_{rg(10)}','user_pw','#CiMPS_2020','allow','printing'])  
    sp.call(['rm',f'./archivos/{filename}-tmp.pdf'])
    #sp.call(['clear'])

def send_email(user,filename,destinatario,Url):
   # Iniciamos los parámetros del script
    remitente ='conferencecimps@cimat.mx'
    
    #destinatario = user['Email']
    #destinatario = 'einar.serna@cimat.mx'
    nombre_adjunto = f'{filename}.pdf'
    ruta_adjunto = f'./archivos/{nombre_adjunto}'

    # Creamos el objeto mensaje
    mensaje = MIMEMultipart()
    
    # Establecemos los atributos del mensaje
    mensaje['From'] = f'CIMPS<{remitente}>'
    mensaje['To'] = destinatario
    mensaje['Subject'] = 'CIMPS 2020: Constancia de Asistencia'
    saludo='Estimad@'
    if user['Genero']=='his':
        saludo='Estimado'
    else:
        saludo='Estimada'

    str_tmp =user['Nombre']
    name= ' '.join(str_tmp[1:])
    dcs=['Mr. ','Alum. ','Dr. ','Miss. ','Prof. ','Comp. ','Eng. ']  
    for dc in dcs:
        if dc in str_tmp:
            str_tmp=str_tmp.split()
            str_tmp= ' '.join(str_tmp[1:])
            name=str_tmp.upper()
        else:
            name=str_tmp.upper()
    # Agregamos el cuerpo del mensaje como objeto MIME de tipo texto
    mensaje.attach(MIMEText(f"""
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Document</title>
</head>
<body>
<p>{saludo} {name}</p>
<p>Agradecemos su participación en la 9na Edición del CIMPS 2020, en el taller de '{user['Taller']}'.</p>
<p>Esperamos contar con su participación para la 10ma, edición del CIMPS 2021, que será llevada a cabo del 20-22 de Octubre del 2021 con sede en la Ciudad de Torreón Coahuila. Si la sindemia del COVID permite realizar actividades presenciales será anunciada a través de las páginas y redes sociales del CIMPS.</p>
<p>No olvide visitar las paginas oficiales: cimps.cimat.mx y facebook/cimps</p>
<p>Adjuntamos su constancia de asistencia, para poder visualizar la constancia y/o realizar una impresión de esta debera utilizar la siguiente contraseña: <strong>#CiMPS_2020</strong></p>
<p>Tambien puede descargar la constancia a través del siguiente <a href="{Url}">enlace</a></p>
<p>Saludos</p>
<p>Comité CIMPS</p>
</body>
</html>    

    """, 'html'))
    
    # Abrimos el archivo que vamos a adjuntar
    
    archivo_adjunto = open(ruta_adjunto, 'rb')
    
    # Creamos un objeto MIME base
    adjunto_MIME = MIMEBase('application', 'octet-stream')
    # Y le cargamos el archivo adjunto
    adjunto_MIME.set_payload((archivo_adjunto).read())
    # Codificamos el objeto en BASE64
    encoders.encode_base64(adjunto_MIME)
    # Agregamos una cabecera al objeto
    adjunto_MIME.add_header('Content-Disposition', 'attachment', filename=nombre_adjunto)
    # Y finalmente lo agregamos al mensaje
    mensaje.attach(adjunto_MIME)
    
    # Creamos la conexión con el servidor
    sesion_smtp = smtplib.SMTP('smtp.gmail.com', 587)
    
    # Ciframos la conexión
    sesion_smtp.starttls()

    # Iniciamos sesión en el servidor
    sesion_smtp.login(remitente,'C1mps2020#')

    # Convertimos el objeto mensaje a texto
    texto = mensaje.as_string()

    # Enviamos el mensaje
    sesion_smtp.sendmail(remitente, destinatario, texto)

    # Cerramos la conexión
    sesion_smtp.quit() 


def main():
    datos=consulta()
    for i,dato in enumerate(datos):
        filename=template(dato)
        pptx2pdf(filename,dato['Email'])
        url=upload_file(filename)
        update_records(dato,filename,url)
        send_email(dato,filename,dato['Email'],url)
        print(f"Correo enviado a : {dato['Email']}\n Enviados: {i+1} de {len(datos)}")
    print('DONE!!')

main()
