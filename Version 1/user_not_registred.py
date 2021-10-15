from credenciales.bd_access import cnxn
from credenciales.google_services import service, MediaFileUpload

from os import getcwd
import subprocess as sp
#library to work with pptx
from pptx import Presentation
#library to work colors of a parragraph in pptx
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

#eventos sin usuarios registrados
def consulta():
    if cnxn!=False:
        with cnxn.cursor() as cursor:
            query="""
                SELECT c.id, c.id_evento, c.id_tipo_constancia, t.texto1, e.contacto, t.texto3, e.nombre_interno 
FROM constancia c
	INNER JOIN tipo_constancia t ON c.id_tipo_constancia = t.id
    INNER JOIN eventos e ON e.id = c.id_evento
WHERE id_tipo_constancia IN (1, 2, 4, 5)
	AND c.id_estado = 1 AND c.id_emision = 9
    AND t.id_estado = 1 AND t.id_emision = 9
	AND c.id > 9999999;
            """
            cursor.execute(query)
            rows = cursor.fetchall()
            data=[]
            for row in rows:
                data.append({
                    'id_constancia':row[0],
                    'id_evento':row[1],
                    'id_tipo_constancia':row[2],
                    'texto1':row[3],
                    'texto2':row[4],
                    'texto3':row[5],
                    'texto4':row[6]
                })
            return data

def update_records(datos,filename,url):
    if cnxn!=False:
        with cnxn.cursor() as cursor:
            query=f"""
                UPDATE 
                    constancia 
                SET 
                    nombre_archivo='{filename}.pdf',
                    url='{url}',
                    generado=1 
                WHERE 
                    id={datos['id_constancia']}
                    AND
                    id_tipo_constancia={datos['id_tipo_constancia']}
            """
            cursor.execute(query)
            cnxn.commit()


def upload_file(filename):
        folder_id = '1nStgZZotbzejT5iOsQj2bY7XjjQnKmhE'
        myfile=f'{filename}.pdf'
        file_metadata = {
            'name': myfile,
            'parents': [folder_id]
        }
        
        media = MediaFileUpload(f'./archivos/{myfile}', mimetype="application/pdf")
        response = service.files().create(
            body=file_metadata,
            media_body=media,
            fields="id"
        ).execute()

        if response:
            url= f"https://drive.google.com/uc?export=download&id={response.get('id')}"
            return url

def template(info):
    #opening template
    document = Presentation('./template/template.pptx')
    #read specific text frame to edit
    #to get access to other frames changes index of shapes and add .text property after .text_frame
    #shape 3 = In recognition ...
    shape_for_recognition = document.slides[0].shapes[3].text_frame
    shape_for_recognition.clear()
    access_for_recognition=shape_for_recognition.paragraphs[0]
    access_for_recognition.alignment = PP_ALIGN.CENTER
    #giving information to paragraph
    texto1=access_for_recognition.add_run()
    texto1.text = info['texto1']
    
    #Giving format to name
    font = texto1.font
    font.name = "Helvetica"
    font.size = Pt(15)
    font.color.rgb = RGBColor(255, 255, 255)




    #Shape 1 Name of author
    shape_for_name = document.slides[0].shapes[1].text_frame
    shape_for_name.clear()
    access_for_name=shape_for_name.paragraphs[0]
    access_for_name.alignment = PP_ALIGN.CENTER
    #giving information to paragraph
    name=access_for_name.add_run()
    name.text = info['texto2']
    
    #Giving format to name
    font = name.font
    font.name = "Helvetica"
    font.size = Pt(20)
    font.color.rgb = RGBColor(255, 191, 0)




    #Shape 4 rol at CIMPS part 1
    shape_for_rol1 = document.slides[0].shapes[4].text_frame
    shape_for_rol1.clear()
    access_for_rol1=shape_for_rol1.paragraphs[0]
    access_for_rol1.alignment = PP_ALIGN.CENTER
    #giving information to paragraph
    texto31=access_for_rol1.add_run()
    texto31.text = f"{info['texto3'].split(',')[0]}"
    
    #Giving format to name
    font = texto31.font
    font.name = "Helvetica"
    font.size = Pt(15)
    font.color.rgb = RGBColor(255, 255, 255)

    #Shape 12 rol at CIMPS part 2
    shape_for_rol2 = document.slides[0].shapes[12].text_frame
    shape_for_rol2.clear()
    access_for_rol2=shape_for_rol2.paragraphs[0]
    access_for_rol2.alignment = PP_ALIGN.CENTER
    #giving information to paragraph
    texto32=access_for_rol2.add_run()
    if len(info['texto3'].split(',')) != 1:
        texto32.text = f"{info['texto3'].split(',')[1]}"
    else:
        texto32.text = ""
    
    #Giving format to name
    font = texto32.font
    font.name = "Helvetica"
    font.size = Pt(15)
    font.color.rgb = RGBColor(255, 255, 255)



    #Shape 10 name of the exposition
    shape_for_expo = document.slides[0].shapes[10].text_frame
    shape_for_expo.clear()
    access_for_expo=shape_for_expo.paragraphs[0]
    access_for_expo.alignment = PP_ALIGN.CENTER
    #giving information to paragraph
    texto4=access_for_expo.add_run()
    texto4.text = info['texto4']
    
    #Giving format to name
    font = texto4.font
    font.name = "Helvetica"
    font.size = Pt(20)
    font.color.rgb = RGBColor(255, 191, 0)

    #Saving file
    filename = f"{info['id_constancia']}_{info['texto4'][0:25].replace(' ','_')}"
    document.save(f"./archivos/{filename}.pptx")
    return filename

def pptx2pdf(filename):
    sp.call(['libreoffice', '--headless', '--convert-to', 'pdf', f'./archivos/{filename}.pptx'])
    sp.call(['mv',f'{filename}.pdf','./archivos'])
    #sp.call(['rm',f'./constancias/{filename}.pptx'])
    #sp.call(['clear'])

def main():
    datos=consulta()
    for dato in datos:
        filename=template(dato)
        pptx2pdf(filename)
        url=upload_file(filename)
        update_records(dato,filename,url)
    print('DONE!!')

main()
