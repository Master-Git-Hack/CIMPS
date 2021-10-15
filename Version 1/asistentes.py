from credenciales.bd_access import cnxn
from credenciales.google_services import service, MediaFileUpload
from random import randrange as rg
from os import getcwd
import subprocess as sp
#library to work with pptx
from pptx import Presentation
#library to work colors of a parragraph in pptx
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

#eventos sin usuarios registrados
def consulta():
    if cnxn!=False:
        with cnxn.cursor() as cursor:
            query="""
SELECT c.id AS Id, t.texto3 as Para, uc.nombre AS Nombre,uc.email as Email, genero as Genero, c.id_tipo_constancia as Tipo_Constancia
FROM constancia c
	INNER JOIN tipo_constancia t ON c.id_tipo_constancia = t.id
    INNER JOIN (
		SELECT id, nombre, email, genero, CASE id_tipo_usuario WHEN 2 THEN 1 WHEN 3 THEN 4 WHEN 4 THEN 5 WHEN 6 THEN 9 END AS id_tipo_constancia
        FROM usuario_constancia
        WHERE id_tipo_usuario = 6
			AND id_estado = 1 AND id_emision = 9
            AND id NOT IN (
				SELECT user_id
				FROM ingsofti_CIMPS3.users_groups
				WHERE group_id = 2
			)
    )uc ON c.id = uc.id AND c.id_tipo_constancia = uc.id_tipo_constancia
WHERE c.id_estado = 1 AND c.id_emision = 9;
            """
            cursor.execute(query)
            rows = cursor.fetchall()
            data=[]
            for row in rows:
#SELECT o2, e.id_locacion, uc.email
                data.append({
                    'Id':row[0],
                    'Para':row[1],
                    'Nombre':row[2],
                    'Email':row[3],
                    'Genero':row[4],
                    'Tipo_Constancia':row[5]
                })
            return data

def update_records(datos,filename,url):
    if cnxn!=False:
        with cnxn.cursor() as cursor:
            query=f"""
                UPDATE 
                    constancia 
                SET 
                    nombre_archivo='{filename}.pdf',
                    url='{url}',
                    generado=1 
                WHERE 
                    id={datos['Id']}
                    AND
                    id_tipo_constancia={datos['Tipo_Constancia']}
            """
            cursor.execute(query)
            cnxn.commit()


def upload_file(filename):
        folder_id = '1TBDKL27T1bDqSrVDX9Ej86iCNJGobjZ7'
        myfile=f'{filename}.pdf'
        file_metadata = {
            'name': myfile,
            'parents': [folder_id]
        }
        
        media = MediaFileUpload(f'./archivos/{myfile}', mimetype="application/pdf")
        response = service.files().create(
            body=file_metadata,
            media_body=media,
            fields="id"
        ).execute()

        if response:
            url= f"https://drive.google.com/uc?export=download&id={response.get('id')}"
            return url

def template(info):
    #opening template
    document = Presentation('./template/template_asistente.pptx')
    #read specific text frame to edit
    #to get access to other frames changes index of shapes and add .text property after .text_frame
    #shape 3 = In recognition ...
    shape_for_recognition = document.slides[0].shapes[4].text_frame
    shape_for_recognition.clear()
    access_for_recognition=shape_for_recognition.paragraphs[0]
    access_for_recognition.alignment = PP_ALIGN.CENTER
    #giving information to paragraph
    texto1=access_for_recognition.add_run()
    
    texto1.text = info['Para'].replace("<hisher>",info['Genero'][0:3])
    
    #Giving format to name
    font = texto1.font
    font.name = "Helvetica"
    font.size = Pt(15)
    font.color.rgb = RGBColor(255, 255, 255)




    #Shape 1 Name of author
    shape_for_name = document.slides[0].shapes[1].text_frame
    shape_for_name.clear()
    access_for_name=shape_for_name.paragraphs[0]
    access_for_name.alignment = PP_ALIGN.CENTER
    #giving information to paragraph
    name=access_for_name.add_run()
    str_tmp =info['Nombre']
    name.text= ' '.join(str_tmp[1:])
    dcs=['Mr. ','Alum. ','Dr. ','Miss. ','Prof. ','Comp. ','Eng. ']  
    for dc in dcs:
        if dc in str_tmp:
            str_tmp=str_tmp.split()
            str_tmp= ' '.join(str_tmp[1:])
            name.text=str_tmp.upper()
        else:
            name.text=str_tmp.upper()
    
    #Giving format to name
    font = name.font
    font.name = "Helvetica"
    font.size = Pt(20)
    font.color.rgb = RGBColor(255, 191, 0)




    #Saving file
    filename = f"{info['Id']}_{info['Email']}"
    document.save(f"./archivos/{filename}-tmp.pptx")
    return filename

def pptx2pdf(filename,email):
    sp.call(['libreoffice', '--headless', '--convert-to', 'pdf', f'./archivos/{filename}-tmp.pptx'])
    #sp.call(['convert',f'./archivos/{filename}.pptx',f'./archivos/{filename}.pdf'])
   # sp.call(['convert',f'./archivos/{filename}.pptx',f'./archivos/{filename}-tmp.pdf'])
    sp.call(['mv',f'{filename}-tmp.pdf','./archivos'])
    sp.call(['pdftk',f'./archivos/{filename}-tmp.pdf','output',f'./archivos/{filename}.pdf','owner_pw',f'#CiMPS_@2020_{email}_{rg(10)}','user_pw','#url','allow','printing'])  
    sp.call(['rm',f'./archivos/{filename}-tmp.pdf'])


def main():
    datos=consulta()
    for dato in datos:
        filename=template(dato)
        pptx2pdf(filename,dato['Email'])
        url=upload_file(filename)
        update_records(dato,filename,url)
    print('DONE!!')

main()